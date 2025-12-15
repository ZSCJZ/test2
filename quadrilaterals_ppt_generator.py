#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
四边形章节PPT生成器
面向8年级学生，教授人教版数学四边形章节
包含平行四边形、特殊的平行四边形和梯形
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np
import os
import tempfile

# 设置matplotlib中文字体
plt.rcParams['font.sans-serif'] = ['SimHei']  # 用来正常显示中文标签
plt.rcParams['axes.unicode_minus'] = False    # 用来正常显示负号

class QuadrilateralsPPTGenerator:
    """四边形PPT生成器类"""
    
    def __init__(self, output_file="四边形.pptx"):
        """初始化PPT生成器"""
        self.prs = Presentation()
        self.output_file = output_file
        self.temp_images = []
    
    def create_cover_slide(self):
        """创建封面幻灯片"""
        slide_layout = self.prs.slide_layouts[0]  # 使用标题幻灯片布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "四边形"
        subtitle.text = "8年级数学 - 人教版"
        
        # 设置标题样式
        title.text_frame.paragraphs[0].font.size = Pt(54)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 78, 152)
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    def create_table_of_contents(self):
        """创建目录幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "目录"
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 78, 152)
        
        # 添加目录项
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        items = [
            "一、四边形的基本概念",
            "二、平行四边形",
            "三、特殊的平行四边形",
            "   3.1 矩形",
            "   3.2 菱形",
            "   3.3 正方形",
            "四、梯形",
            "五、总结与练习"
        ]
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            
            # 根据缩进级别设置不同的项目符号
            if item.startswith("  "):
                p.level = 1
            else:
                p.level = 0
    
    def _add_picture_to_slide(self, slide, img_path, width=Cm(8)):
        """统一的图片添加方法，确保图片位置合理，不与文本重叠"""
        # 将图片放在右侧，距离左侧14cm，顶部6cm，避免与文本区域重叠
        left = Cm(14)
        top = Cm(6)
        slide.shapes.add_picture(img_path, left, top, width=width)
    
    def create_basic_concepts_slide(self):
        """创建四边形基本概念幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "四边形的基本概念"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 78, 152)
        
        # 添加平行四边形定义和基本说明
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        
        # 添加定义段落
        p1 = tf.add_paragraph()
        p1.text = "定义：由不在同一直线上的四条线段首尾顺次相接组成的封闭图形叫做四边形。"
        p1.font.size = Pt(18)
        p1.line_spacing = 1.5
        
        # 添加构成要素
        p2 = tf.add_paragraph()
        p2.text = "四边形的构成："
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.space_after = Pt(12)
        
        # 构成要素列表
        elements = [
            "顶点：四边形的四个端点",
            "边：连接顶点的四条线段",
            "内角：四边形内部的四个角",
            "对角线：连接不相邻顶点的线段"
        ]
        
        for element in elements:
            p = tf.add_paragraph()
            p.text = element
            p.font.size = Pt(16)
            p.level = 1
            
        # 绘制一个简单的四边形图示
        fig, ax = self._create_figure()
        
        # 绘制四边形
        x = [0, 2, 3, 1]
        y = [0, 0, 2, 2]
        ax.plot(x + [x[0]], y + [y[0]], 'b-', linewidth=2)
        
        # 标记顶点
        labels = ['A', 'B', 'C', 'D']
        for i, (xi, yi, label) in enumerate(zip(x, y, labels)):
            ax.annotate(label, (xi, yi), fontsize=14, 
                        xytext=(10, 10), textcoords='offset points')
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-0.5, 3.5)
        ax.set_ylim(-0.5, 2.5)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def create_parallelogram_intro(self):
        """创建平行四边形介绍幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "平行四边形"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)
        
        # 添加定义和基本说明
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        # 定义段落
        p1 = tf.add_paragraph()
        p1.text = "定义：两组对边分别平行的四边形叫做平行四边形。"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.space_after = Pt(12)
        
        # 平行四边形特点
        p2 = tf.add_paragraph()
        p2.text = "平行四边形表示方法："
        p2.font.size = Pt(16)
        p2.space_after = Pt(6)
        
        p3 = tf.add_paragraph()
        p3.text = "用符号□表示，例如：□ABCD"
        p3.font.size = Pt(16)
        p3.level = 1
        
        # 绘制平行四边形
        fig, ax = self._create_figure()
        
        # 绘制平行四边形
        x = [0, 3, 4, 1]
        y = [0, 0, 2, 2]
        ax.plot(x + [x[0]], y + [y[0]], 'g-', linewidth=2)
        
        # 标记顶点
        labels = ['A', 'B', 'C', 'D']
        for i, (xi, yi, label) in enumerate(zip(x, y, labels)):
            ax.annotate(label, (xi, yi), fontsize=14, 
                        xytext=(10, 10), textcoords='offset points')
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-0.5, 4.5)
        ax.set_ylim(-0.5, 2.5)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def create_parallelogram_properties(self):
        """创建平行四边形性质幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "平行四边形的性质"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)
        
        # 添加性质列表
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        properties = [
            "1. 对边平行且相等",
            "2. 对角相等",
            "3. 邻角互补",
            "4. 对角线互相平分",
            "5. 是中心对称图形，对称中心是对角线的交点"
        ]
        
        for prop in properties:
            p = tf.add_paragraph()
            p.text = prop
            p.font.size = Pt(18)
            p.space_after = Pt(6)
        
        # 绘制带有性质标注的平行四边形
        fig, ax = self._create_figure()
        
        # 绘制平行四边形
        x = [0, 4, 5, 1]
        y = [0, 0, 3, 3]
        ax.plot(x + [x[0]], y + [y[0]], 'g-', linewidth=2)
        
        # 标记顶点
        labels = ['A', 'B', 'C', 'D']
        for i, (xi, yi, label) in enumerate(zip(x, y, labels)):
            ax.annotate(label, (xi, yi), fontsize=14, 
                        xytext=(10, 10), textcoords='offset points')
        
        # 绘制对角线
        ax.plot([x[0], x[2]], [y[0], y[2]], 'r--', linewidth=1)  # AC对角线
        ax.plot([x[1], x[3]], [y[1], y[3]], 'r--', linewidth=1)  # BD对角线
        
        # 标记对角线交点
        intersection = [(x[0]+x[2])/2, (y[0]+y[2])/2]
        ax.plot(intersection[0], intersection[1], 'ro', markersize=6)
        ax.annotate('O', intersection, fontsize=12, 
                   xytext=(10, -10), textcoords='offset points')
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-1, 6)
        ax.set_ylim(-1, 4)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def create_parallelogram_theorems(self):
        """创建平行四边形判定定理幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "平行四边形的判定"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)
        
        # 添加判定定理
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        theorems = [
            "1. 两组对边分别平行的四边形是平行四边形（定义）",
            "2. 两组对边分别相等的四边形是平行四边形",
            "3. 一组对边平行且相等的四边形是平行四边形",
            "4. 对角线互相平分的四边形是平行四边形",
            "5. 两组对角分别相等的四边形是平行四边形"
        ]
        
        for theorem in theorems:
            p = tf.add_paragraph()
            p.text = theorem
            p.font.size = Pt(16)
            p.space_after = Pt(6)
    
    def create_rectangle_slide(self):
        """创建矩形幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "特殊的平行四边形 - 矩形"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 140, 0)
        
        # 添加矩形定义和基本性质
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        # 定义段落
        p1 = tf.add_paragraph()
        p1.text = "定义：有一个角是直角的平行四边形叫做矩形（长方形）。"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.space_after = Pt(12)
        
        # 矩形性质
        p2 = tf.add_paragraph()
        p2.text = "矩形的性质："
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.space_after = Pt(6)
        
        properties = [
            "1. 具有平行四边形的所有性质",
            "2. 四个角都是直角",
            "3. 对角线相等且互相平分",
            "4. 既是中心对称图形，又是轴对称图形"
        ]
        
        for prop in properties:
            p = tf.add_paragraph()
            p.text = prop
            p.font.size = Pt(16)
            p.level = 1
        
        # 绘制矩形
        fig, ax = self._create_figure()
        
        # 绘制矩形（直角平行四边形）
        x = [0, 4, 4, 0]
        y = [0, 0, 3, 3]
        ax.plot(x + [x[0]], y + [y[0]], 'orange', linewidth=2)
        
        # 标记直角
        ax.plot([0.8, 0, 0], [0, 0, 0.8], 'orange', linewidth=1)
        
        # 标记顶点
        labels = ['A', 'B', 'C', 'D']
        for i, (xi, yi, label) in enumerate(zip(x, y, labels)):
            ax.annotate(label, (xi, yi), fontsize=14, 
                        xytext=(10, 10), textcoords='offset points')
        
        # 绘制对角线
        ax.plot([x[0], x[2]], [y[0], y[2]], 'r--', linewidth=1)  # AC对角线
        ax.plot([x[1], x[3]], [y[1], y[3]], 'r--', linewidth=1)  # BD对角线
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-0.5, 4.5)
        ax.set_ylim(-0.5, 3.5)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def create_rhombus_slide(self):
        """创建菱形幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "菱形"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(218, 165, 32)
        
        # 添加菱形定义和基本性质
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        # 定义段落
        p1 = tf.add_paragraph()
        p1.text = "定义：有一组邻边相等的平行四边形叫做菱形。"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.space_after = Pt(12)
        
        # 菱形性质
        p2 = tf.add_paragraph()
        p2.text = "菱形的性质："
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.space_after = Pt(6)
        
        properties = [
            "1. 具有平行四边形的所有性质",
            "2. 四条边都相等",
            "3. 对角线互相垂直且平分",
            "4. 对角线平分一组对角",
            "5. 既是中心对称图形，又是轴对称图形"
        ]
        
        for prop in properties:
            p = tf.add_paragraph()
            p.text = prop
            p.font.size = Pt(16)
            p.level = 1
        
        # 绘制菱形
        fig, ax = self._create_figure()
        
        # 绘制菱形（邻边相等的平行四边形）
        x = [1, 3, 1, -1]
        y = [0, 2, 4, 2]
        ax.plot(x + [x[0]], y + [y[0]], 'purple', linewidth=2)
        
        # 标记顶点
        labels = ['A', 'B', 'C', 'D']
        for i, (xi, yi, label) in enumerate(zip(x, y, labels)):
            ax.annotate(label, (xi, yi), fontsize=14, 
                        xytext=(10, 10), textcoords='offset points')
        
        # 绘制对角线
        ax.plot([x[0], x[2]], [y[0], y[2]], 'r--', linewidth=1)  # AC对角线
        ax.plot([x[1], x[3]], [y[1], y[3]], 'r--', linewidth=1)  # BD对角线
        
        # 标记垂直符号
        mid_x, mid_y = (x[0]+x[2])/2, (y[0]+y[2])/2
        ax.plot([mid_x-0.2, mid_x, mid_x+0.2], [mid_y-0.2, mid_y, mid_y+0.2], 'r-', linewidth=1)
        ax.plot([mid_x+0.2, mid_x, mid_x-0.2], [mid_y-0.2, mid_y, mid_y+0.2], 'r-', linewidth=1)
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-1.5, 3.5)
        ax.set_ylim(-0.5, 4.5)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def create_square_slide(self):
        """创建正方形幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "特殊的平行四边形 - 正方形"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 0, 0)
        
        # 添加内容
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        # 定义段落
        p1 = tf.add_paragraph()
        p1.text = "定义：有一组邻边相等并且有一个角是直角的平行四边形叫做正方形。"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.space_after = Pt(12)
        
        # 正方形性质
        p2 = tf.add_paragraph()
        p2.text = "正方形的性质："
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.space_after = Pt(6)
        
        properties = [
            "1. 具有平行四边形、矩形、菱形的所有性质",
            "2. 四条边都相等",
            "3. 四个角都是直角",
            "4. 对角线相等且互相垂直平分",
            "5. 对角线平分一组对角",
            "6. 既是中心对称图形，又是轴对称图形"
        ]
        
        for prop in properties:
            p = tf.add_paragraph()
            p.text = prop
            p.font.size = Pt(16)
            p.level = 1
        
        # 绘制正方形
        fig, ax = self._create_figure()
        
        # 绘制正方形（四边相等且有直角的平行四边形）
        x = [0, 3, 3, 0]
        y = [0, 0, 3, 3]
        ax.plot(x + [x[0]], y + [y[0]], 'red', linewidth=2)
        
        # 标记直角
        ax.plot([0.8, 0, 0], [0, 0, 0.8], 'red', linewidth=1)
        
        # 标记顶点
        labels = ['A', 'B', 'C', 'D']
        for i, (xi, yi, label) in enumerate(zip(x, y, labels)):
            ax.annotate(label, (xi, yi), fontsize=14, 
                        xytext=(10, 10), textcoords='offset points')
        
        # 绘制对角线
        ax.plot([x[0], x[2]], [y[0], y[2]], 'r--', linewidth=1)  # AC对角线
        ax.plot([x[1], x[3]], [y[1], y[3]], 'r--', linewidth=1)  # BD对角线
        
        # 标记垂直符号
        mid_x, mid_y = 1.5, 1.5
        ax.plot([mid_x-0.2, mid_x, mid_x+0.2], [mid_y-0.2, mid_y, mid_y+0.2], 'r-', linewidth=1)
        ax.plot([mid_x+0.2, mid_x, mid_x-0.2], [mid_y-0.2, mid_y, mid_y+0.2], 'r-', linewidth=1)
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-0.5, 3.5)
        ax.set_ylim(-0.5, 3.5)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def create_special_parallelogram_relationship(self):
        """创建特殊平行四边形关系幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "特殊平行四边形之间的关系"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(72, 61, 139)
        
        # 添加内容
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        relationships = [
            "1. 矩形、菱形、正方形都是特殊的平行四边形",
            "2. 正方形既是矩形，又是菱形",
            "3. 矩形和菱形不一定是正方形",
            "4. 平行四边形不一定是矩形、菱形或正方形"
        ]
        
        for rel in relationships:
            p = tf.add_paragraph()
            p.text = rel
            p.font.size = Pt(18)
            p.space_after = Pt(6)
    
    def create_trapezoid_intro(self):
        """创建梯形定义幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "梯形的定义"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(46, 139, 87)
        
        # 添加内容
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        # 定义段落
        p1 = tf.add_paragraph()
        p1.text = "定义：一组对边平行，另一组对边不平行的四边形叫做梯形。"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.space_after = Pt(12)
        
        # 梯形各部分名称
        p2 = tf.add_paragraph()
        p2.text = "梯形的各部分名称："
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.space_after = Pt(6)
        
        parts = [
            "1. 平行的两边叫做梯形的底边（上底和下底）",
            "2. 不平行的两边叫做梯形的腰",
            "3. 两腰中点的连线叫做梯形的中位线",
            "4. 梯形的高：从一底上的任一点向另一底作垂线，这点和垂足之间的线段叫做梯形的高"
        ]
        
        for part in parts:
            p = tf.add_paragraph()
            p.text = part
            p.font.size = Pt(16)
            p.level = 1
        
        # 绘制普通梯形
        fig, ax = self._create_figure()
        
        # 绘制梯形
        x = [0, 4, 3, 1]
        y = [0, 0, 3, 3]
        ax.plot(x + [x[0]], y + [y[0]], 'green', linewidth=2)
        
        # 标记平行符号
        ax.plot([-0.3, 0, 0], [0, 0, -0.3], 'green', linewidth=1)
        ax.plot([2.7, 3, 3], [3, 3, 3.3], 'green', linewidth=1)
        
        # 标记顶点
        labels = ['A', 'B', 'C', 'D']
        for i, (xi, yi, label) in enumerate(zip(x, y, labels)):
            ax.annotate(label, (xi, yi), fontsize=14, 
                        xytext=(10, 10), textcoords='offset points')
        
        # 绘制高
        ax.plot([1, 1], [0, 3], '--b', linewidth=1)  # 高
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-0.5, 4.5)
        ax.set_ylim(-0.5, 3.5)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def create_trapezoid_classification(self):
        """创建梯形分类幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "梯形的分类"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(46, 139, 87)
        
        # 添加梯形分类内容
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        # 分类介绍
        p1 = tf.add_paragraph()
        p1.text = "梯形可分为以下几类："
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.space_after = Pt(6)
        
        # 分类列表
        classifications = [
            "1. 一般梯形：两腰不相等的梯形",
            "2. 等腰梯形：两腰相等的梯形",
            "3. 直角梯形：有一个角是直角的梯形"
        ]
        
        for cls in classifications:
            p = tf.add_paragraph()
            p.text = cls
            p.font.size = Pt(16)
            p.level = 1
        
        # 绘制等腰梯形
        fig, ax = self._create_figure()
        
        # 绘制等腰梯形（两腰相等）
        x = [0, 4, 3, 1]
        y = [0, 0, 3, 3]
        ax.plot(x + [x[0]], y + [y[0]], 'green', linewidth=2)
        
        # 标记平行符号
        ax.plot([-0.3, 0, 0], [0, 0, -0.3], 'green', linewidth=1)
        ax.plot([2.7, 3, 3], [3, 3, 3.3], 'green', linewidth=1)
        
        # 标记顶点
        labels = ['A', 'B', 'C', 'D']
        for i, (xi, yi, label) in enumerate(zip(x, y, labels)):
            ax.annotate(label, (xi, yi), fontsize=14, 
                        xytext=(10, 10), textcoords='offset points')
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-0.5, 4.5)
        ax.set_ylim(-0.5, 3.5)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def create_trapezoid_properties(self):
        """创建梯形性质幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "梯形的性质"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(46, 139, 87)
        
        # 添加梯形性质内容
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        # 梯形一般性质
        p1 = tf.add_paragraph()
        p1.text = "一般梯形的性质："
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.space_after = Pt(6)
        
        general_props = [
            "1. 梯形的中位线平行于两底",
            "2. 梯形的中位线长度等于两底和的一半",
            "3. 梯形的面积等于（上底+下底）× 高 ÷ 2"
        ]
        
        for prop in general_props:
            p = tf.add_paragraph()
            p.text = prop
            p.font.size = Pt(16)
            p.level = 1
        
        # 等腰梯形性质
        p2 = tf.add_paragraph()
        p2.text = "等腰梯形的性质："
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.space_after = Pt(6)
        
        isosceles_props = [
            "1. 两腰相等",
            "2. 同一底上的两个角相等",
            "3. 对角线相等",
            "4. 是轴对称图形，对称轴是上下底中点的连线"
        ]
        
        for prop in isosceles_props:
            p = tf.add_paragraph()
            p.text = prop
            p.font.size = Pt(16)
            p.level = 1
        
        # 绘制直角梯形
        fig, ax = self._create_figure()
        
        # 绘制直角梯形（有一个角是直角）
        x = [0, 4, 4, 0]
        y = [0, 0, 3, 3]
        ax.plot(x + [x[0]], y + [y[0]], 'green', linewidth=2)
        
        # 标记直角
        ax.plot([0.8, 0, 0], [0, 0, 0.8], 'green', linewidth=1)
        ax.plot([3.2, 4, 4], [0, 0, 0.8], 'green', linewidth=1)
        
        # 标记顶点
        labels = ['A', 'B', 'C', 'D']
        for i, (xi, yi, label) in enumerate(zip(x, y, labels)):
            ax.annotate(label, (xi, yi), fontsize=14, 
                        xytext=(10, 10), textcoords='offset points')
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-0.5, 4.5)
        ax.set_ylim(-0.5, 3.5)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def create_summary_slide(self):
        """创建总结幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "四边形知识总结"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 0, 0)
        
        # 添加总结内容
        content = slide.placeholders[1]
        # 限制文本区域宽度，避免与图片重叠
        content.width = Cm(12)
        tf = content.text_frame
        tf.clear()
        
        # 总结要点
        summary_points = [
            "1. 四边形的基本概念：由不在同一直线上的四条线段首尾顺次连接而成的图形",
            "2. 平行四边形：两组对边分别平行的四边形，具有对边相等、对角相等、对角线互相平分等性质",
            "3. 特殊平行四边形：",
            "   - 矩形：有一个角是直角的平行四边形，具有四个直角、对角线相等的性质",
            "   - 菱形：有一组邻边相等的平行四边形，具有四边相等、对角线互相垂直的性质",
            "   - 正方形：既是矩形又是菱形，具有矩形和菱形的所有性质",
            "4. 梯形：一组对边平行，另一组对边不平行的四边形",
            "   - 等腰梯形：两腰相等，同一底上的两个角相等",
            "   - 直角梯形：有一个角是直角的梯形"
        ]
        
        for point in summary_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.space_after = Pt(6)
    
    def create_exercises_slide(self):
        """创建练习题幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "练习题"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
        
        # 添加内容
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # 练习题
        exercises = [
            "1. 平行四边形ABCD中，AB=6cm，BC=8cm，求平行四边形的周长。",
            "2. 矩形的一条对角线长为10cm，一边长为6cm，求另一边长。",
            "3. 菱形的对角线分别为8cm和6cm，求菱形的面积和边长。",
            "4. 梯形的上底为5cm，下底为10cm，高为4cm，求梯形的面积。",
            "5. 正方形的周长为24cm，求正方形的面积。"
        ]
        
        for i, exercise in enumerate(exercises, 1):
            p = tf.add_paragraph()
            p.text = exercise
            p.font.size = Pt(16)
            p.space_after = Pt(12)
        
        # 绘制一个包含四边形关系的示意图
        fig, ax = self._create_figure()
        
        # 绘制四边形的包含关系
        ax.add_patch(plt.Rectangle((-2, -2), 8, 8, fill=False, edgecolor='black', linewidth=2))
        ax.text(2, 5, '四边形', fontsize=16, ha='center')
        
        # 平行四边形
        ax.add_patch(plt.Rectangle((-1, -1), 3, 4, fill=False, edgecolor='green', linewidth=2))
        ax.text(0.5, 3, '平行四边形', fontsize=14, ha='center', color='green')
        
        # 梯形
        ax.add_patch(plt.Polygon([[3, -1], [6, -1], [5, 1], [2, 3]], fill=False, edgecolor='orange', linewidth=2))
        ax.text(4, 0.5, '梯形', fontsize=14, ha='center', color='orange')
        
        # 矩形
        ax.add_patch(plt.Rectangle((-0.5, 0.5), 2, 2, fill=False, edgecolor='blue', linewidth=2))
        ax.text(0.5, 2, '矩形', fontsize=12, ha='center', color='blue')
        
        # 菱形
        ax.add_patch(plt.Polygon([[0, -0.5], [1, 0.5], [0, 1.5], [-1, 0.5]], fill=False, edgecolor='purple', linewidth=2))
        ax.text(0, 0.5, '菱形', fontsize=12, ha='center', color='purple')
        
        # 正方形
        ax.add_patch(plt.Rectangle([0, 0.5], 1, 1, fill=False, edgecolor='red', linewidth=2))
        ax.text(0.5, 1.2, '正方形', fontsize=10, ha='center', color='red')
        
        # 隐藏坐标轴
        ax.axis('off')
        
        # 设置图形范围
        ax.set_xlim(-2.5, 6.5)
        ax.set_ylim(-2.5, 6.5)
        
        # 保存临时图片
        img_path = self._save_temp_image(fig)
        
        # 使用统一方法在幻灯片中插入图片
        self._add_picture_to_slide(slide, img_path)
    
    def save(self):
        """保存PPT文件"""
        self.prs.save(self.output_file)
        print(f"PPT已保存到: {self.output_file}")
        
        # 清理临时图片文件
        for img_path in self.temp_images:
            try:
                os.remove(img_path)
            except:
                pass
    
    def _create_figure(self, title=None):
        """创建一个matplotlib图形"""
        fig, ax = plt.subplots(figsize=(8, 6))
        ax.set_aspect('equal')
        if title:
            ax.set_title(title, fontsize=16)
        return fig, ax
    
    def _save_temp_image(self, fig):
        """保存临时图片并返回路径"""
        fd, path = tempfile.mkstemp(suffix='.png')
        os.close(fd)
        fig.savefig(path, dpi=300, bbox_inches='tight')
        plt.close(fig)
        self.temp_images.append(path)
        return path

# 主函数
def main():
    """主函数"""
    # 创建PPT生成器实例
    ppt = QuadrilateralsPPTGenerator("四边形.pptx")
    
    print("开始生成四边形PPT...")
    
    # 创建基本框架
    ppt.create_cover_slide()
    ppt.create_table_of_contents()
    ppt.create_basic_concepts_slide()
    
    # 创建平行四边形章节
    print("创建平行四边形章节...")
    ppt.create_parallelogram_intro()
    ppt.create_parallelogram_properties()
    ppt.create_parallelogram_theorems()
    
    # 创建特殊平行四边形章节
    print("创建特殊平行四边形章节...")
    ppt.create_special_parallelogram_relationship()  # 先介绍关系
    ppt.create_rectangle_slide()  # 矩形
    ppt.create_rhombus_slide()    # 菱形
    ppt.create_square_slide()     # 正方形
    
    # 创建梯形章节
    print("创建梯形章节...")
    ppt.create_trapezoid_intro()       # 梯形定义
    ppt.create_trapezoid_classification()  # 梯形分类
    ppt.create_trapezoid_properties()   # 梯形性质
    
    # 创建总结和练习题
    print("创建总结和练习题...")
    ppt.create_summary_slide()         # 总结幻灯片
    ppt.create_exercises_slide()       # 练习题幻灯片
    
    # 保存PPT文件
    ppt.save()
    
if __name__ == "__main__":
    main()